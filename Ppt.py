import streamlit as st
import pandas as pd
import io
import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from langchain_groq import ChatGroq
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_classic.chains import LLMChain
from langchain_classic.prompts import PromptTemplate
import tempfile
import re

# Streamlit UI Configuration
st.set_page_config(
    page_title="AutoPPT Pro - AI Presentation Generator",
    page_icon="📊",
    layout="wide"
)

# Title and Description
st.title("🤖 AutoPPT Pro - AI Presentation Generator")
st.markdown("""
Generate professional PowerPoint presentations automatically using Generative AI. 
Upload your data or describe your topic, and let AI create a complete presentation for you.
""")

# API Key Section in Sidebar
st.sidebar.header("🔑 API Configuration")

# Groq API Key input
groq_api_key = st.sidebar.text_input(
    "Enter your Groq API Key:",
    type="password",
    placeholder="gsk_xxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxx",
    help="Get your free API key from https://console.groq.com"
)

# Check if API key is provided
if not groq_api_key:
    st.sidebar.warning("⚠️ Please enter your Groq API key to continue")
    st.stop()

# Initialize LLM with user-provided API key
try:
    llm = ChatGroq(groq_api_key=groq_api_key, model_name="openai/gpt-oss-120b")
    # Test the API key with a simple call
    llm.invoke("Hello")  # Simple test to verify API key
    st.sidebar.success("✅ Groq API key validated!")
except Exception as e:
    st.sidebar.error(f"❌ Invalid Groq API key: {str(e)}")
    st.stop()

# Initialize Hugging Face embeddings
try:
    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2",
        model_kwargs={'device': 'cpu'}
    )
    st.sidebar.success("✅ Hugging Face embeddings loaded!")
except Exception as e:
    st.sidebar.error(f"❌ Error loading embeddings: {str(e)}")
    st.stop()

# Main Content Area
tab1, tab2, tab3 = st.tabs(["📝 Text to PPT", "📊 Data to PPT", "⚙️ Advanced Settings"])

def create_presentation_from_structure(topic, audience, slides_count, tone, additional_context):
    """Create PowerPoint presentation using AI-generated structure"""
    
    # Create presentation structure using LangChain with Groq
    ppt_prompt = PromptTemplate(
        input_variables=["topic", "audience", "slides_count", "tone", "additional_context"],
        template="""
        Create a comprehensive PowerPoint presentation structure for:
        
        TOPIC: {topic}
        AUDIENCE: {audience}
        NUMBER OF SLIDES: {slides_count}
        TONE: {tone}
        ADDITIONAL CONTEXT: {additional_context}
        
        Please create exactly {slides_count} slides with this structure:
        
        SLIDE 1: TITLE SLIDE
        - Main Title: {topic}
        - Subtitle: Presentation for {audience}
        
        SLIDE 2: AGENDA
        - Key point 1
        - Key point 2  
        - Key point 3
        
        SLIDE 3-{last_slide_minus_one}: CONTENT SLIDES
        - Each with a clear heading
        - 3-5 bullet points per slide
        - Focus on key information for {audience}
        
        SLIDE {slides_count}: CONCLUSION
        - Summary of main points
        - Key takeaways
        - Next steps or recommendations
        
        Provide the content in this exact format:
        SLIDE 1: TITLE | {topic} | Presentation for {audience}
        SLIDE 2: AGENDA | Overview of Topics | Key Areas Covered | Main Discussion Points
        SLIDE 3: [Heading 1] | [Bullet 1] | [Bullet 2] | [Bullet 3]
        SLIDE 4: [Heading 2] | [Bullet 1] | [Bullet 2] | [Bullet 3]
        ...
        SLIDE {slides_count}: CONCLUSION | Summary | Key Takeaways | Recommendations
        """
    )
    
    ppt_chain = LLMChain(llm=llm, prompt=ppt_prompt)
    ppt_structure = ppt_chain.run({
        "topic": topic,
        "audience": audience,
        "slides_count": slides_count,
        "tone": tone,
        "additional_context": additional_context,
        "last_slide_minus_one": slides_count - 1
    })
    
    return ppt_structure

def parse_slide_structure(ppt_structure):
    """Parse the AI-generated structure into slide components - FIXED VERSION"""
    slides = []
    
    # Split by lines and look for SLIDE patterns
    lines = ppt_structure.strip().split('\n')
    
    for line in lines:
        line = line.strip()
        # Look for lines that start with "SLIDE X:" pattern
        if re.match(r'^SLIDE\s+\d+:', line, re.IGNORECASE) or re.match(r'^Slide\s+\d+:', line, re.IGNORECASE):
            # Extract slide number and content
            slide_match = re.match(r'^(?:SLIDE|Slide)\s+(\d+):\s*(.+)', line, re.IGNORECASE)
            if slide_match:
                slide_content = slide_match.group(2).strip()
                slides.append(slide_content)
        
        # Also check for numbered slides without "SLIDE" prefix
        elif re.match(r'^\d+\.', line) or re.match(r'^\d+\)', line):
            slide_content = re.sub(r'^\d+[\.\)]\s*', '', line).strip()
            if slide_content:
                slides.append(slide_content)
    
    # If no structured slides found, try to extract any meaningful content
    if not slides:
        st.warning("⚠️ Using fallback parsing method...")
        # Try to split by common slide indicators
        potential_slides = re.split(r'\n\s*\n|\d+\.\s|\d+\)\s|SLIDE\s+\d+:', ppt_structure)
        for potential_slide in potential_slides:
            slide_content = potential_slide.strip()
            if slide_content and len(slide_content) > 20:  # Only take substantial content
                slides.append(slide_content)
    
    return slides

def create_ppt_file(slides_data, topic, audience):
    """Create actual PowerPoint file from slide data - IMPROVED VERSION"""
    prs = Presentation()
    
    # If no slides were parsed, create a default presentation
    if not slides_data:
        st.warning("⚠️ No slides parsed. Creating default presentation structure.")
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = topic
        subtitle.text = f"Presentation for {audience}"
        
        # Add some default slides
        default_slides = [
            ("Agenda", ["Introduction", "Key Points", "Conclusion"]),
            ("Introduction", ["Topic overview", "Purpose of presentation", "Target audience"]),
            ("Key Points", ["Main point 1", "Main point 2", "Main point 3"]),
            ("Conclusion", ["Summary", "Key takeaways", "Next steps"])
        ]
        
        for slide_title, bullet_points in default_slides:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            
            title_shape.text = slide_title
            tf = body_shape.text_frame
            tf.text = bullet_points[0]
            
            for point in bullet_points[1:]:
                p = tf.add_paragraph()
                p.text = point
    else:
        # Process parsed slides
        for i, slide_content in enumerate(slides_data):
            if i == 0:
                # First slide as title slide
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                # Try to extract title and subtitle from content
                if '|' in slide_content:
                    parts = [part.strip() for part in slide_content.split('|')]
                    title.text = parts[0] if len(parts) > 0 else topic
                    subtitle.text = parts[1] if len(parts) > 1 else f"Presentation for {audience}"
                else:
                    title.text = slide_content[:50]  # First 50 chars as title
                    subtitle.text = f"Presentation for {audience}"
            else:
                # Content slides
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                title_shape = slide.shapes.title
                body_shape = slide.placeholders[1]
                
                # Parse slide content
                if '|' in slide_content:
                    parts = [part.strip() for part in slide_content.split('|')]
                    title_text = parts[0] if parts else f"Slide {i+1}"
                    bullet_points = parts[1:] if len(parts) > 1 else ["Content to be added"]
                else:
                    title_text = f"Slide {i+1}"
                    bullet_points = [slide_content] if slide_content else ["Content to be added"]
                
                title_shape.text = title_text
                
                # Add bullet points
                if bullet_points:
                    tf = body_shape.text_frame
                    tf.text = bullet_points[0] if bullet_points else "Key points"
                    
                    for point in bullet_points[1:]:
                        if point and point != "CONTENT":
                            p = tf.add_paragraph()
                            p.text = point
    
    # Save to bytes
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    
    return ppt_bytes, len(prs.slides)

with tab1:
    st.header("Generate PPT from Text Description")
    
    topic = st.text_input("Presentation Topic:", placeholder="e.g., Digital Transformation Strategy 2024")
    audience = st.selectbox("Target Audience:", ["Executive Leadership", "Technical Team", "Sales & Marketing", "General Audience", "Students"])
    slides_count = st.slider("Number of Slides:", 5, 20, 10)
    tone = st.selectbox("Presentation Tone:", ["Professional", "Persuasive", "Educational", "Inspirational", "Technical"])
    
    additional_context = st.text_area("Additional Context (Optional):", 
                                    placeholder="Any specific points to include, company information, or special requirements...")
    
    if st.button("Generate Presentation", key="text_to_ppt"):
        if not topic:
            st.error("Please enter a presentation topic!")
        else:
            with st.spinner("🤖 AI is generating your presentation..."):
                try:
                    # Generate presentation structure
                    ppt_structure = create_presentation_from_structure(
                        topic, audience, slides_count, tone, additional_context
                    )
                    
                    # Parse the structure
                    slides_data = parse_slide_structure(ppt_structure)
                    
                    # Create PowerPoint file
                    ppt_bytes, actual_slides_count = create_ppt_file(slides_data, topic, audience)
                    
                    # Download button
                    st.success("🎉 Presentation generated successfully!")
                    
                    col1, col2 = st.columns([1, 2])
                    
                    with col1:
                        st.download_button(
                            label="📥 Download PowerPoint",
                            data=ppt_bytes,
                            file_name=f"{topic.replace(' ', '_')}_presentation.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    
                    with col2:
                        st.info(f"**Generated {actual_slides_count} slides** using Groq LLM and Hugging Face embeddings")
                    
                    # Show generated content
                    with st.expander("📋 View AI-Generated Content"):
                        st.write("### Raw AI Response")
                        st.text_area("Structure:", ppt_structure, height=200)
                        
                        st.write("### Parsed Slides")
                        if slides_data:
                            for i, slide in enumerate(slides_data):
                                st.write(f"**Slide {i+1}:** {slide}")
                        else:
                            st.warning("No slides could be parsed from the AI response. Using default structure.")
                            
                except Exception as e:
                    st.error(f"Error generating presentation: {str(e)}")
                    import traceback
                    st.code(traceback.format_exc())

# ... (rest of the code remains the same for tab2 and tab3)

with tab2:
    st.header("Generate PPT from Data")
    
    uploaded_file = st.file_uploader("Upload CSV/Excel file", type=['csv', 'xlsx'])
    
    if uploaded_file:
        # Save uploaded file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix=".csv" if uploaded_file.name.endswith('.csv') else ".xlsx") as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_file_path = tmp_file.name
        
        try:
            if uploaded_file.name.endswith('.csv'):
                df = pd.read_csv(tmp_file_path)
            else:
                df = pd.read_excel(tmp_file_path)
            
            st.write("### Data Preview")
            st.dataframe(df.head())
            
            analysis_type = st.selectbox("Analysis Type:", 
                                       ["Data Summary", "Trend Analysis", "Comparative Analysis", "Key Insights"])
            
            if st.button("Generate Data Presentation", key="data_to_ppt"):
                with st.spinner("🤖 Analyzing data and creating presentation..."):
                    try:
                        # Simple data presentation without complex parsing
                        prs = Presentation()
                        
                        # Title slide
                        title_slide_layout = prs.slide_layouts[0]
                        slide = prs.slides.add_slide(title_slide_layout)
                        title = slide.shapes.title
                        subtitle = slide.placeholders[1]
                        title.text = f"Data Analysis: {analysis_type}"
                        subtitle.text = f"Based on {uploaded_file.name}"
                        
                        # Data overview slide
                        bullet_slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(bullet_slide_layout)
                        title_shape = slide.shapes.title
                        body_shape = slide.placeholders[1]
                        title_shape.text = "Dataset Overview"
                        tf = body_shape.text_frame
                        tf.text = f"Total Rows: {len(df)}"
                        p = tf.add_paragraph()
                        p.text = f"Total Columns: {len(df.columns)}"
                        p = tf.add_paragraph()
                        p.text = f"Columns: {', '.join(df.columns.tolist())}"
                        
                        # Summary statistics slide
                        slide = prs.slides.add_slide(bullet_slide_layout)
                        title_shape = slide.shapes.title
                        body_shape = slide.placeholders[1]
                        title_shape.text = "Summary Statistics"
                        tf = body_shape.text_frame
                        
                        numeric_cols = df.select_dtypes(include=['number']).columns
                        if len(numeric_cols) > 0:
                            tf.text = f"Numeric columns: {len(numeric_cols)}"
                            for col in numeric_cols[:3]:  # Show first 3 numeric columns
                                p = tf.add_paragraph()
                                p.text = f"{col}: mean={df[col].mean():.2f}"
                        else:
                            tf.text = "No numeric columns found"
                        
                        # Save presentation
                        ppt_bytes = io.BytesIO()
                        prs.save(ppt_bytes)
                        ppt_bytes.seek(0)
                        
                        st.success("🎉 Data presentation generated!")
                        st.download_button(
                            label="📥 Download Data Presentation",
                            data=ppt_bytes,
                            file_name=f"data_analysis_{analysis_type.replace(' ', '_')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                        
                    except Exception as e:
                        st.error(f"Error in data analysis: {str(e)}")
        
        except Exception as e:
            st.error(f"Error reading file: {str(e)}")
        finally:
            # Clean up temporary file
            import os
            if os.path.exists(tmp_file_path):
                os.unlink(tmp_file_path)

with tab3:
    st.header("Advanced Settings")
    
    st.subheader("Presentation Template")
    template_choice = st.selectbox("Choose Template:", 
                                 ["Corporate Blue", "Modern Red", "Professional Green", "Creative Purple"])
    
    st.subheader("AI Model Settings")
    st.info(f"Current Model: Groq - openai/gpt-oss-120b")
    st.info(f"Embeddings: Hugging Face - sentence-transformers/all-MiniLM-L6-v2")
    
    st.subheader("Content Options")
    col1, col2 = st.columns(2)
    
    with col1:
        include_agenda = st.checkbox("Include Agenda Slide", value=True)
        include_summary = st.checkbox("Include Summary Slide", value=True)
        include_qa = st.checkbox("Include Q&A Slide", value=True)
    
    with col2:
        add_speaker_notes = st.checkbox("Add Speaker Notes", value=False)
        add_references = st.checkbox("Add References Slide", value=False)
        detailed_bullets = st.checkbox("Detailed Bullet Points", value=True)

# Instructions Section
with st.expander("📖 How to Use This Tool"):
    st.markdown("""
    ### Step-by-Step Guide:
    
    1. **Configure API**: Enter your Groq API key in the sidebar
    2. **Choose Input Method**:
       - **Text to PPT**: Describe your topic and requirements
       - **Data to PPT**: Upload your dataset for analysis-based presentations
    3. **Customize Settings**: Adjust audience, tone, and slide count
    4. **Generate**: Click the generate button and wait for AI to create your presentation
    5. **Download**: Get your professionally formatted PowerPoint file
    
    ### Technical Stack:
    - **LLM**: Groq (openai/gpt-oss-120b) for content generation
    - **Embeddings**: Hugging Face (sentence-transformers/all-MiniLM-L6-v2)
    - **Framework**: LangChain for orchestration
    - **UI**: Streamlit for web interface
    - **PPT Generation**: python-pptx library
    
    ### Features:
    - AI-powered content generation using Groq LLM
    - Professional PowerPoint formatting
    - Data analysis and visualization recommendations
    - Multiple template options
    - Enterprise-ready presentation output
    """)

# Footer
st.markdown("---")
st.markdown(
    "🔒 Built with Groq LLM & Hugging Face Embeddings | " +
    "Powered by LangChain & Streamlit",
    unsafe_allow_html=True
)
